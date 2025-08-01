"""
Report Generation System for Multifamily Property Analysis
Generates professional reports in Excel, PDF, and PowerPoint formats
"""

import os
import json
import pandas as pd
from typing import Dict, List, Optional, Any
import logging
from datetime import datetime
from pathlib import Path

try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.chart import LineChart, BarChart, Reference
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logging.warning("openpyxl not available. Excel generation disabled.")

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("ReportLab not available. PDF generation disabled.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint generation disabled.")

class ReportGenerator:
    """Generates professional reports in multiple formats"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.logger = logging.getLogger(__name__)
        
        # Report styling configuration
        self.brand_colors = {
            'primary': '#1f4e79',
            'secondary': '#70ad47', 
            'accent': '#ffc000',
            'text': '#333333',
            'light_gray': '#f2f2f2'
        }
    
    def generate_reports(self,
                        documents: Dict[str, Any],
                        ai_insights: Dict[str, Any],
                        financial_analysis: Dict[str, Any],
                        output_dir: str,
                        job_id: str,
                        generate_pitch_deck: bool = True,
                        include_analysis: bool = True) -> List[Dict[str, str]]:
        """
        Generate all requested reports
        
        Args:
            documents: Parsed document data
            ai_insights: AI analysis results  
            financial_analysis: Financial models and projections
            output_dir: Directory to save reports
            job_id: Unique job identifier
            generate_pitch_deck: Whether to generate PowerPoint pitch deck
            include_analysis: Whether to include detailed analysis Excel file
            
        Returns:
            List of generated report files with metadata
        """
        
        self.logger.info(f"Generating reports for job {job_id}")
        
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        generated_reports = []
        
        # Generate integrated data JSON (always generated)
        json_report = self._generate_json_report(
            documents, ai_insights, financial_analysis, output_dir, job_id
        )
        if json_report:
            generated_reports.append(json_report)
        
        # Generate Excel analysis report
        if include_analysis and EXCEL_AVAILABLE:
            excel_report = self._generate_excel_report(
                documents, ai_insights, financial_analysis, output_dir, job_id
            )
            if excel_report:
                generated_reports.append(excel_report)
        
        # Generate PDF executive summary
        if PDF_AVAILABLE:
            pdf_report = self._generate_pdf_report(
                documents, ai_insights, financial_analysis, output_dir, job_id
            )
            if pdf_report:
                generated_reports.append(pdf_report)
        
        # Generate PowerPoint pitch deck
        if generate_pitch_deck and PPTX_AVAILABLE:
            pptx_report = self._generate_powerpoint_report(
                documents, ai_insights, financial_analysis, output_dir, job_id
            )
            if pptx_report:
                generated_reports.append(pptx_report)
        
        self.logger.info(f"Generated {len(generated_reports)} reports")
        return generated_reports
    
    def _generate_json_report(self, documents: Dict[str, Any], ai_insights: Dict[str, Any], 
                            financial_analysis: Dict[str, Any], output_dir: str, job_id: str) -> Optional[Dict[str, str]]:
        """Generate comprehensive JSON data report"""
        
        try:
            # Compile comprehensive data
            integrated_data = {
                'job_id': job_id,
                'generated_at': datetime.now().isoformat(),
                'property_summary': self._extract_property_summary(documents, ai_insights),
                'financial_summary': self._extract_financial_summary(financial_analysis),
                'ai_insights_summary': self._extract_ai_summary(ai_insights),
                'raw_data': {
                    'documents': documents,
                    'ai_insights': ai_insights,
                    'financial_analysis': financial_analysis
                }
            }
            
            # Save JSON file
            json_path = os.path.join(output_dir, 'integratedData.json')
            with open(json_path, 'w') as f:
                json.dump(integrated_data, f, indent=2, default=str)
            
            return {
                'type': 'json',
                'filename': 'integratedData.json',
                'path': json_path,
                'description': 'Comprehensive property analysis data'
            }
            
        except Exception as e:
            self.logger.error(f"Error generating JSON report: {e}")
            return None
    
    def _generate_excel_report(self, documents: Dict[str, Any], ai_insights: Dict[str, Any],
                             financial_analysis: Dict[str, Any], output_dir: str, job_id: str) -> Optional[Dict[str, str]]:
        """Generate detailed Excel analysis report"""
        
        if not EXCEL_AVAILABLE:
            return None
        
        try:
            # Create workbook
            wb = Workbook()
            
            # Remove default sheet
            wb.remove(wb.active)
            
            # Create sheets
            self._create_executive_summary_sheet(wb, documents, ai_insights, financial_analysis)
            self._create_financial_projections_sheet(wb, financial_analysis)
            self._create_unit_analysis_sheet(wb, documents)
            self._create_returns_analysis_sheet(wb, financial_analysis)
            self._create_sensitivity_analysis_sheet(wb, financial_analysis)
            
            # Save workbook
            excel_path = os.path.join(output_dir, 'populatedTemplate.xlsx')
            wb.save(excel_path)
            
            return {
                'type': 'excel',
                'filename': 'populatedTemplate.xlsx', 
                'path': excel_path,
                'description': 'Detailed financial analysis workbook'
            }
            
        except Exception as e:
            self.logger.error(f"Error generating Excel report: {e}")
            return None
    
    def _create_executive_summary_sheet(self, wb: Workbook, documents: Dict[str, Any], 
                                      ai_insights: Dict[str, Any], financial_analysis: Dict[str, Any]):
        """Create executive summary sheet"""
        
        ws = wb.create_sheet("Executive Summary")
        
        # Title
        ws['A1'] = "Multifamily Property Analysis - Executive Summary"
        ws['A1'].font = Font(size=16, bold=True)
        
        # Property Information
        row = 3
        ws[f'A{row}'] = "Property Information"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 1
        
        property_info = self._extract_property_summary(documents, ai_insights)
        for key, value in property_info.items():
            ws[f'A{row}'] = key.replace('_', ' ').title()
            ws[f'B{row}'] = str(value)
            row += 1
        
        # Financial Summary
        row += 2
        ws[f'A{row}'] = "Financial Summary"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 1
        
        financial_summary = self._extract_financial_summary(financial_analysis)
        for key, value in financial_summary.items():
            ws[f'A{row}'] = key.replace('_', ' ').title()
            ws[f'B{row}'] = value if isinstance(value, str) else f"${value:,.0f}"
            row += 1
        
        # AI Insights
        row += 2
        ws[f'A{row}'] = "Key Insights"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 1
        
        ai_summary = self._extract_ai_summary(ai_insights)
        for insight in ai_summary.get('key_insights', []):
            ws[f'A{row}'] = f"• {insight}"
            row += 1
    
    def _create_financial_projections_sheet(self, wb: Workbook, financial_analysis: Dict[str, Any]):
        """Create financial projections sheet"""
        
        ws = wb.create_sheet("Financial Projections")
        
        projections = financial_analysis.get('investment_projections', {}).get('yearly_projections', [])
        
        if not projections:
            ws['A1'] = "No projection data available"
            return
        
        # Headers
        headers = ['Year', 'Gross Income', 'Effective Gross Income', 'Operating Expenses', 'NOI', 'NOI Growth %']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col)
            cell.value = header
            cell.font = Font(bold=True)
        
        # Data
        for row, projection in enumerate(projections, 2):
            ws.cell(row=row, column=1).value = projection.get('year', '')
            ws.cell(row=row, column=2).value = projection.get('gross_income', 0)
            ws.cell(row=row, column=3).value = projection.get('effective_gross_income', 0)
            ws.cell(row=row, column=4).value = projection.get('total_expenses', 0)
            ws.cell(row=row, column=5).value = projection.get('noi', 0)
            ws.cell(row=row, column=6).value = projection.get('noi_growth', 0)
        
        # Add chart
        self._add_projection_chart(ws, len(projections))
    
    def _create_unit_analysis_sheet(self, wb: Workbook, documents: Dict[str, Any]):
        """Create unit analysis sheet"""
        
        ws = wb.create_sheet("Unit Analysis")
        
        rent_roll = documents.get('rent_roll', {})
        units = rent_roll.get('units', [])
        
        if not units:
            ws['A1'] = "No unit data available"
            return
        
        # Headers
        headers = ['Unit', 'Bedrooms', 'Bathrooms', 'Sq Ft', 'Current Rent', 'Rent/Sq Ft', 'Status']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col)
            cell.value = header
            cell.font = Font(bold=True)
        
        # Data
        for row, unit in enumerate(units, 2):
            ws.cell(row=row, column=1).value = unit.get('unit', '')
            ws.cell(row=row, column=2).value = unit.get('bedrooms', 0)
            ws.cell(row=row, column=3).value = unit.get('bathrooms', 0)
            ws.cell(row=row, column=4).value = unit.get('sqft', 0)
            ws.cell(row=row, column=5).value = unit.get('current_rent', 0)
            
            # Calculate rent per sq ft
            sqft = unit.get('sqft', 0)
            rent = unit.get('current_rent', 0)
            rent_per_sqft = rent / sqft if sqft > 0 else 0
            ws.cell(row=row, column=6).value = round(rent_per_sqft, 2)
            
            ws.cell(row=row, column=7).value = unit.get('status', 'Unknown')
    
    def _create_returns_analysis_sheet(self, wb: Workbook, financial_analysis: Dict[str, Any]):
        """Create returns analysis sheet"""
        
        ws = wb.create_sheet("Returns Analysis")
        
        returns = financial_analysis.get('returns_analysis', {})
        
        # Key metrics
        row = 1
        ws[f'A{row}'] = "Investment Returns Analysis"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 2
        
        metrics = {
            'Estimated Current Value': returns.get('estimated_value', 0),
            'Exit Value': returns.get('exit_value', 0),
            'Value Appreciation': returns.get('value_appreciation', 0),
            'Current Cap Rate': f"{returns.get('cap_rate', 0):.2f}%",
            'Exit Cap Rate': f"{returns.get('exit_cap_rate', 0):.2f}%",
            'NOI Growth (Total)': f"{returns.get('noi_growth_total', 0):.2f}%"
        }
        
        for metric, value in metrics.items():
            ws[f'A{row}'] = metric
            ws[f'B{row}'] = value
            row += 1
    
    def _create_sensitivity_analysis_sheet(self, wb: Workbook, financial_analysis: Dict[str, Any]):
        """Create sensitivity analysis sheet"""
        
        ws = wb.create_sheet("Sensitivity Analysis")
        
        sensitivity = financial_analysis.get('sensitivity_analysis', {})
        
        # Rent growth sensitivity
        row = 1
        ws[f'A{row}'] = "Rent Growth Sensitivity"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 2
        
        ws[f'A{row}'] = "Rent Growth %"
        ws[f'B{row}'] = "Exit Value"
        row += 1
        
        for scenario in sensitivity.get('rent_growth_sensitivity', []):
            ws[f'A{row}'] = f"{scenario.get('rent_growth', 0)}%"
            ws[f'B{row}'] = scenario.get('exit_value', 0)
            row += 1
    
    def _add_projection_chart(self, ws, num_years: int):
        """Add projection chart to worksheet"""
        
        try:
            # Create chart
            chart = LineChart()
            chart.title = "NOI Projections"
            chart.x_axis.title = "Year"
            chart.y_axis.title = "NOI ($)"
            
            # Data reference
            data = Reference(ws, min_col=5, min_row=1, max_row=num_years + 1)
            categories = Reference(ws, min_col=1, min_row=2, max_row=num_years + 1)
            
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(categories)
            
            # Add chart to worksheet
            ws.add_chart(chart, "H2")
            
        except Exception as e:
            self.logger.warning(f"Could not add chart: {e}")
    
    def _generate_pdf_report(self, documents: Dict[str, Any], ai_insights: Dict[str, Any],
                           financial_analysis: Dict[str, Any], output_dir: str, job_id: str) -> Optional[Dict[str, str]]:
        """Generate PDF executive summary report"""
        
        if not PDF_AVAILABLE:
            return None
        
        try:
            pdf_path = os.path.join(output_dir, 'analysisReport.pdf')
            
            # Create PDF document
            doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
                alignment=1  # Center alignment
            )
            
            # Title
            story.append(Paragraph("Multifamily Property Analysis Report", title_style))
            story.append(Spacer(1, 20))
            
            # Property summary
            property_summary = self._extract_property_summary(documents, ai_insights)
            story.append(Paragraph("Property Summary", styles['Heading2']))
            
            summary_data = []
            for key, value in property_summary.items():
                summary_data.append([key.replace('_', ' ').title(), str(value)])
            
            if summary_data:
                summary_table = Table(summary_data, colWidths=[2.5*inch, 3*inch])
                summary_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(summary_table)
            
            story.append(Spacer(1, 20))
            
            # Financial summary
            story.append(Paragraph("Financial Analysis", styles['Heading2']))
            financial_summary = self._extract_financial_summary(financial_analysis)
            
            financial_data = []
            for key, value in financial_summary.items():
                formatted_value = value if isinstance(value, str) else f"${value:,.0f}"
                financial_data.append([key.replace('_', ' ').title(), formatted_value])
            
            if financial_data:
                financial_table = Table(financial_data, colWidths=[2.5*inch, 3*inch])
                financial_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(financial_table)
            
            story.append(PageBreak())
            
            # AI insights
            story.append(Paragraph("Key Insights & Recommendations", styles['Heading2']))
            ai_summary = self._extract_ai_summary(ai_insights)
            
            for insight in ai_summary.get('key_insights', []):
                story.append(Paragraph(f"• {insight}", styles['Normal']))
                story.append(Spacer(1, 10))
            
            # Build PDF
            doc.build(story)
            
            return {
                'type': 'pdf',
                'filename': 'analysisReport.pdf',
                'path': pdf_path,
                'description': 'Executive summary report'
            }
            
        except Exception as e:
            self.logger.error(f"Error generating PDF report: {e}")
            return None
    
    def _generate_powerpoint_report(self, documents: Dict[str, Any], ai_insights: Dict[str, Any],
                                  financial_analysis: Dict[str, Any], output_dir: str, job_id: str) -> Optional[Dict[str, str]]:
        """Generate PowerPoint pitch deck"""
        
        if not PPTX_AVAILABLE:
            return None
        
        try:
            # Create presentation
            prs = Presentation()
            
            # Title slide
            self._create_title_slide(prs, documents, job_id)
            
            # Property overview slide
            self._create_property_overview_slide(prs, documents, ai_insights)
            
            # Financial summary slide
            self._create_financial_summary_slide(prs, financial_analysis)
            
            # Investment highlights slide
            self._create_investment_highlights_slide(prs, ai_insights)
            
            # Returns analysis slide
            self._create_returns_slide(prs, financial_analysis)
            
            # Market analysis slide
            self._create_market_analysis_slide(prs, ai_insights)
            
            # Next steps slide
            self._create_next_steps_slide(prs)
            
            # Save presentation
            pptx_path = os.path.join(output_dir, 'pitchDeck.pptx')
            prs.save(pptx_path)
            
            return {
                'type': 'powerpoint',
                'filename': 'pitchDeck.pptx',
                'path': pptx_path,
                'description': 'Professional investor pitch deck'
            }
            
        except Exception as e:
            self.logger.error(f"Error generating PowerPoint report: {e}")
            return None
    
    def _create_title_slide(self, prs: Presentation, documents: Dict[str, Any], job_id: str):
        """Create title slide"""
        
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        property_info = self._extract_property_summary(documents, {})
        property_name = property_info.get('property_name', 'Multifamily Property')
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{property_name} Investment Opportunity"
        subtitle.text = f"Property Analysis Report\\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
    
    def _create_property_overview_slide(self, prs: Presentation, documents: Dict[str, Any], ai_insights: Dict[str, Any]):
        """Create property overview slide"""
        
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Property Overview"
        
        content = slide.placeholders[1]
        
        property_summary = self._extract_property_summary(documents, ai_insights)
        
        overview_text = ""
        for key, value in property_summary.items():
            overview_text += f"{key.replace('_', ' ').title()}: {value}\\n"
        
        content.text = overview_text
    
    def _create_financial_summary_slide(self, prs: Presentation, financial_analysis: Dict[str, Any]):
        """Create financial summary slide"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Financial Summary"
        
        content = slide.placeholders[1]
        
        financial_summary = self._extract_financial_summary(financial_analysis)
        
        summary_text = ""
        for key, value in financial_summary.items():
            formatted_value = value if isinstance(value, str) else f"${value:,.0f}"
            summary_text += f"{key.replace('_', ' ').title()}: {formatted_value}\\n"
        
        content.text = summary_text
    
    def _create_investment_highlights_slide(self, prs: Presentation, ai_insights: Dict[str, Any]):
        """Create investment highlights slide"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Investment Highlights"
        
        content = slide.placeholders[1]
        
        ai_summary = self._extract_ai_summary(ai_insights)
        highlights = ai_summary.get('key_insights', ['Strong property fundamentals', 'Good market location', 'Value-add potential'])
        
        highlights_text = ""
        for highlight in highlights[:5]:  # Limit to 5 highlights
            highlights_text += f"• {highlight}\\n"
        
        content.text = highlights_text
    
    def _create_returns_slide(self, prs: Presentation, financial_analysis: Dict[str, Any]):
        """Create returns analysis slide"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Projected Returns"
        
        content = slide.placeholders[1]
        
        returns = financial_analysis.get('returns_analysis', {})
        
        returns_text = f"""Cap Rate: {returns.get('cap_rate', 0):.2f}%
Exit Cap Rate: {returns.get('exit_cap_rate', 0):.2f}%
Estimated Value: ${returns.get('estimated_value', 0):,.0f}
Exit Value: ${returns.get('exit_value', 0):,.0f}
Total Appreciation: ${returns.get('value_appreciation', 0):,.0f}"""
        
        content.text = returns_text
    
    def _create_market_analysis_slide(self, prs: Presentation, ai_insights: Dict[str, Any]):
        """Create market analysis slide"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Market Analysis"
        
        content = slide.placeholders[1]
        
        market_analysis = ai_insights.get('market_analysis', {})
        
        market_text = f"""Market: {market_analysis.get('market', 'Strong multifamily market')}
Occupancy Analysis: {ai_insights.get('property_analysis', {}).get('occupancy_analysis', {}).get('occupancy_status', 'Good')}
Investment Grade: {ai_insights.get('summary', {}).get('investment_grade', 'B')}"""
        
        content.text = market_text
    
    def _create_next_steps_slide(self, prs: Presentation):
        """Create next steps slide"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps"
        
        content = slide.placeholders[1]
        
        next_steps = """• Review detailed financial analysis
• Conduct property inspection
• Verify market assumptions
• Finalize financing terms
• Submit letter of intent"""
        
        content.text = next_steps
    
    # Helper methods to extract summaries
    
    def _extract_property_summary(self, documents: Dict[str, Any], ai_insights: Dict[str, Any]) -> Dict[str, Any]:
        """Extract property summary information"""
        
        summary = {}
        
        # From documents
        if documents.get('rent_roll'):
            rent_roll = documents['rent_roll']
            summary.update({
                'total_units': len(rent_roll.get('units', [])),
                'total_monthly_rent': rent_roll.get('summary', {}).get('total_monthly_rent', 0),
                'average_rent': rent_roll.get('summary', {}).get('average_rent', 0)
            })
        
        if documents.get('offering_memo'):
            property_info = documents['offering_memo'].get('property_info', {})
            summary.update(property_info)
        
        # From AI insights
        property_analysis = ai_insights.get('property_analysis', {})
        occupancy_analysis = property_analysis.get('occupancy_analysis', {})
        
        if occupancy_analysis:
            summary['occupancy_rate'] = f"{occupancy_analysis.get('occupancy_rate', 0):.1f}%"
        
        return summary
    
    def _extract_financial_summary(self, financial_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Extract financial summary information"""
        
        summary = {}
        
        base_financials = financial_analysis.get('base_financials', {}).get('current_financials', {})
        
        summary.update({
            'gross_income': base_financials.get('gross_income', 0),
            'operating_expenses': base_financials.get('operating_expenses', 0),
            'noi': base_financials.get('noi', 0),
            'expense_ratio': f"{base_financials.get('expense_ratio', 0):.1f}%"
        })
        
        returns = financial_analysis.get('returns_analysis', {})
        if returns:
            summary.update({
                'estimated_value': returns.get('estimated_value', 0),
                'cap_rate': f"{returns.get('cap_rate', 0):.2f}%"
            })
        
        return summary
    
    def _extract_ai_summary(self, ai_insights: Dict[str, Any]) -> Dict[str, Any]:
        """Extract AI insights summary"""
        
        summary = {'key_insights': []}
        
        # Extract key insights from various analysis sections
        property_analysis = ai_insights.get('property_analysis', {})
        financial_insights = ai_insights.get('financial_insights', {})
        investment_summary = ai_insights.get('summary', {})
        
        # Add occupancy insights
        occupancy = property_analysis.get('occupancy_analysis', {})
        if occupancy.get('occupancy_rate'):
            summary['key_insights'].append(f"Property maintains {occupancy['occupancy_rate']:.1f}% occupancy rate")
        
        # Add revenue insights
        revenue_ops = property_analysis.get('revenue_optimization', {})
        if revenue_ops.get('monthly_revenue_potential', 0) > 0:
            summary['key_insights'].append(f"${revenue_ops['monthly_revenue_potential']:,.0f} monthly revenue upside identified")
        
        # Add investment grade
        if investment_summary.get('investment_grade'):
            summary['key_insights'].append(f"Investment grade: {investment_summary['investment_grade']}")
        
        # Add strengths
        for strength in investment_summary.get('key_strengths', [])[:2]:
            summary['key_insights'].append(strength)
        
        return summary